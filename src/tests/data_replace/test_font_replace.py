import os

from pptx.dml.color import RGBColor

from .testcase import PPTTTestCase
from ...PPTT.ppt import make_ppt

BASE_DIR = os.path.dirname(os.path.abspath(__file__))


class TextFontReplaceTestCase(PPTTTestCase):
    # debug = True
    # temp_dir = BASE_DIR
    master_slide = os.path.join(BASE_DIR, 'text_font_replace.pptx')

    def test_replace_font_bold(self):
        ms = self.get_master_ppt()
        slide_pos = 2
        real_slide_pos = slide_pos - 1
        origin_slide = ms.slides[real_slide_pos]
        # check placeholder exists
        style_shape_name = 'bold_text'
        plain_shape_name = 'plain'
        origin_bold_shape = self.find_shape(origin_slide, style_shape_name)
        self.assertTrue(origin_bold_shape)
        self.assertTrue(origin_bold_shape.text_frame.paragraphs[0].runs[0].font.bold)
        self.assertTrue(self.find_shape(origin_slide, plain_shape_name))

        input = [
            {
                "slide_pos": slide_pos,
                "contents": {
                    style_shape_name: {
                        "text": {
                            "font": {
                                "bold": False
                            },
                            "value": "text is not bold"
                        },
                    },
                    plain_shape_name: {
                        "text": {
                            "font": {
                                "bold": True
                            },
                            "value": "now text in bold",
                        },
                    }
                }
            }
        ]
        make_ppt(self.master_slide, self.target_slide, input)
        ppt = self.get_target_ppt()
        slide = ppt.slides[real_slide_pos]

        style_shape = self.find_shape(slide, style_shape_name)
        style_font = style_shape.text_frame.paragraphs[0].runs[0].font

        self.assertFalse(style_font.bold)

        plain_shape = self.find_shape(slide, plain_shape_name)
        plain_font = plain_shape.text_frame.paragraphs[0].runs[0].font

        self.assertTrue(plain_font.bold)

    def test_replace_font_italic(self):
        ms = self.get_master_ppt()
        slide_pos = 2
        real_slide_pos = slide_pos - 1
        origin_slide = ms.slides[real_slide_pos]
        # check placeholder exists
        style_shape_name = 'italic_text'
        plain_shape_name = 'plain'
        origin_bold_shape = self.find_shape(origin_slide, style_shape_name)
        self.assertTrue(origin_bold_shape)
        self.assertTrue(origin_bold_shape.text_frame.paragraphs[0].runs[0].font.italic)
        self.assertTrue(self.find_shape(origin_slide, plain_shape_name))

        input = [
            {
                "slide_pos": slide_pos,
                "contents": {
                    style_shape_name: {
                        "text": {
                            "font": {
                                "italic": False
                            },
                            "value": "text is not bold"
                        },
                    },
                    plain_shape_name: {
                        "text": {
                            "font": {
                                "italic": True
                            },
                            "value": "now text in bold",
                        },
                    }
                }
            }
        ]
        make_ppt(self.master_slide, self.target_slide, input)
        ppt = self.get_target_ppt()
        slide = ppt.slides[real_slide_pos]

        style_shape = self.find_shape(slide, style_shape_name)
        style_font = style_shape.text_frame.paragraphs[0].runs[0].font

        self.assertFalse(style_font.italic)

        plain_shape = self.find_shape(slide, plain_shape_name)
        plain_font = plain_shape.text_frame.paragraphs[0].runs[0].font

        self.assertTrue(plain_font.italic)

    def test_replace_font_underline(self):
        ms = self.get_master_ppt()
        slide_pos = 2
        real_slide_pos = slide_pos - 1
        origin_slide = ms.slides[real_slide_pos]
        # check placeholder exists
        style_shape_name = 'underline_text'
        plain_shape_name = 'plain'
        origin_bold_shape = self.find_shape(origin_slide, style_shape_name)
        self.assertTrue(origin_bold_shape)
        self.assertTrue(origin_bold_shape.text_frame.paragraphs[0].runs[0].font.underline)
        self.assertTrue(self.find_shape(origin_slide, plain_shape_name))

        input = [
            {
                "slide_pos": slide_pos,
                "contents": {
                    style_shape_name: {
                        "text": {
                            "font": {
                                "underline": False
                            },
                            "value": "text is not bold"
                        },
                    },
                    plain_shape_name: {
                        "text": {
                            "font": {
                                "underline": True
                            },
                            "value": "now text in bold",
                        },
                    }
                }
            }
        ]
        make_ppt(self.master_slide, self.target_slide, input)
        ppt = self.get_target_ppt()
        slide = ppt.slides[real_slide_pos]

        style_shape = self.find_shape(slide, style_shape_name)
        style_font = style_shape.text_frame.paragraphs[0].runs[0].font

        self.assertFalse(style_font.underline)

        plain_shape = self.find_shape(slide, plain_shape_name)
        plain_font = plain_shape.text_frame.paragraphs[0].runs[0].font

        self.assertTrue(plain_font.underline)

    def test_replace_font_color(self):
        ms = self.get_master_ppt()
        slide_pos = 1
        real_slide_pos = slide_pos - 1
        origin_slide = ms.slides[real_slide_pos]
        # check placeholder exists
        title_shape_name = 'title'
        body_shape_name = 'body'
        self.assertTrue(self.find_shape(origin_slide, title_shape_name))
        self.assertTrue(self.find_shape(origin_slide, body_shape_name))

        input = [
            {
                "slide_pos": slide_pos,
                "contents": {
                    title_shape_name: {
                        "text": {
                            "font": {
                                "color": "#219653",
                            },
                            "value": "this is test title"
                        },
                    },
                    body_shape_name: {
                        "text": {
                            "font": {
                                "color": "#EB5757",
                            },
                            "value": "test \nbody",
                        },
                    }
                }
            }
        ]
        make_ppt(self.master_slide, self.target_slide, input)
        ppt = self.get_target_ppt()
        slide = ppt.slides[real_slide_pos]
        self.assertEqual(2, len(slide.shapes))

        contents = input[0]['contents']
        title_shape = self.find_shape(slide, title_shape_name)
        title_font = title_shape.text_frame.paragraphs[0].font
        font_color = title_font.color.rgb

        title_hex_color = contents[title_shape_name]['text']['font']['color'][1:]
        self.assertEqual(RGBColor.from_string(title_hex_color), font_color)

        body_shape = self.find_shape(slide, body_shape_name)
        body_font = body_shape.text_frame.paragraphs[0].runs[0].font
        body_color = body_font.color.rgb

        body_hex_color = contents[body_shape_name]['text']['font']['color'][1:]
        self.assertEqual(RGBColor.from_string(body_hex_color), body_color)


class TableFontReplaceTestCase(PPTTTestCase):
    # debug = True
    # temp_dir = BASE_DIR

    master_slide = os.path.join(BASE_DIR, 'text_font_replace.pptx')

    def test_key_value_type(self):
        ms = self.get_master_ppt()
        slide_pos = 3
        real_slide_pos = slide_pos - 1
        origin_slide = ms.slides[real_slide_pos]
        # check placeholder exists
        shape_name = 'table'
        self.assertTrue(self.find_shape(origin_slide, shape_name))
        black_color = '000000'
        green_color = '00F900'
        input = [
            {
                "slide_pos": slide_pos,
                "contents": {
                    shape_name: {
                        "table": {
                            "data_type": "key_value",
                            "keys": [
                                {
                                    "name": {"value": "now black", "font": {"color": f"#{black_color}"}},
                                    "data_key": "color"
                                },
                                {
                                    "name": {"value": "no bold", "font": {"bold": False}},
                                    "data_key": "bold"
                                },
                                {
                                    "name": {"value": "no underline", "font": {"underline": False}},
                                    "data_key": "underline"
                                },
                                {
                                    "name": {"value": "no italic", "font": {"italic": False}},
                                    "data_key": "italic"
                                },
                            ],
                            "data": [
                                {
                                    "color": {"value": "this is green", "font": {"color": f"#{green_color}"}},
                                    "bold": {"value": "this is bold", "font": {"bold": True}},
                                    "italic": {"value": "this is italic", "font": {"italic": True}},
                                    "underline": {"value": "this is underline", "font": {"underline": True}},
                                },
                            ]
                        }
                    },
                }
            }
        ]

        make_ppt(self.master_slide, self.target_slide, input)
        ppt = self.get_target_ppt()
        slide = ppt.slides[real_slide_pos]

        table_shape = self.find_shape(slide, shape_name)
        color_text, bold_text, underline_text, italic_text = (cell.text_frame.paragraphs[0].runs[0].font for
                                                              cell in table_shape.table.rows[0].cells)
        self.assertEqual(RGBColor.from_string(black_color), color_text.color.rgb)
        self.assertFalse(bold_text.bold)
        self.assertFalse(underline_text.underline)
        self.assertFalse(italic_text.italic)

        color_text, bold_text, underline_text, italic_text = (cell.text_frame.paragraphs[0].runs[0].font for
                                                              cell in table_shape.table.rows[1].cells)
        self.assertEqual(RGBColor.from_string(green_color), color_text.color.rgb)
        self.assertTrue(bold_text.bold)
        self.assertTrue(underline_text.underline)
        self.assertTrue(italic_text.italic)

    def test_raw_type(self):
        ms = self.get_master_ppt()
        slide_pos = 3
        real_slide_pos = slide_pos - 1
        origin_slide = ms.slides[real_slide_pos]
        # check placeholder exists
        shape_name = 'table'
        self.assertTrue(self.find_shape(origin_slide, shape_name))
        black_color = '000000'
        green_color = '00F900'
        input = [
            {
                "slide_pos": slide_pos,
                "contents": {
                    shape_name: {
                        "table": {
                            "data_type": "raw",
                            "data": [
                                [
                                    {"value": "now black", "font": {"color": f"#{black_color}"}},
                                    {"value": "no bold", "font": {"bold": False}},
                                    {"value": "no underline", "font": {"underline": False}},
                                    {"value": "no italic", "font": {"italic": False}},
                                ],
                                [
                                    {"value": "this is green", "font": {"color": f"#{green_color}"}},
                                    {"value": "this is bold", "font": {"bold": True}},
                                    {"value": "this is underline", "font": {"underline": True}},
                                    {"value": "this is italic", "font": {"italic": True}},
                                ],

                            ],
                        }
                    },
                }
            }
        ]
        make_ppt(self.master_slide, self.target_slide, input)
        ppt = self.get_target_ppt()
        slide = ppt.slides[real_slide_pos]

        table_shape = self.find_shape(slide, shape_name)
        color_text, bold_text, underline_text, italic_text = (cell.text_frame.paragraphs[0].runs[0].font for
                                                              cell in table_shape.table.rows[0].cells)
        self.assertEqual(RGBColor.from_string(black_color), color_text.color.rgb)
        self.assertFalse(bold_text.bold)
        self.assertFalse(underline_text.underline)
        self.assertFalse(italic_text.italic)

        color_text, bold_text, underline_text, italic_text = (cell.text_frame.paragraphs[0].runs[0].font for
                                                              cell in table_shape.table.rows[1].cells)
        self.assertEqual(RGBColor.from_string(green_color), color_text.color.rgb)
        self.assertTrue(bold_text.bold)
        self.assertTrue(underline_text.underline)
        self.assertTrue(italic_text.italic)
