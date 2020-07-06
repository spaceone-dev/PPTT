import os
import tempfile
from unittest import TestCase

from pptx import Presentation
from pptx.presentation import Presentation as PresentationType

from ...PPTT.ppt import make_ppt, find_shape

BASE_DIR = os.path.dirname(os.path.abspath(__file__))


class PPTTTestCase(TestCase):
    master_slide = ''
    temp_dir = None
    target_slide = ''

    def setUp(self):
        self.temp_dir = tempfile.TemporaryDirectory()
        self.target_slide = os.path.join(self.temp_dir.name, 'target.pptx')

    def tearDown(self):
        self.temp_dir.cleanup()
        self.temp_dir = None

    def get_master_ppt(self) -> PresentationType:
        return Presentation(self.master_slide)

    def get_target_ppt(self) -> PresentationType:
        return Presentation(self.target_slide)

    @staticmethod
    def find_shape(slide, name):
        return find_shape(slide, name, soft=True)


class TextPlaceholderTestCase(PPTTTestCase):
    master_slide = os.path.join(BASE_DIR, 'text_placeholder.pptx')

    def test_replace_text(self):
        ms = self.get_master_ppt()
        slide_layout_name = 'simple_title'
        origin_layout = ms.slide_master.slide_layouts.get_by_name(slide_layout_name)
        # check placeholder exists
        title_shape_name = 'title'
        body_shape_name = 'body'
        self.assertTrue([s for s in origin_layout.shapes if s.name == title_shape_name][0])
        self.assertTrue([s for s in origin_layout.shapes if s.name == body_shape_name][0])

        input = [
            {
                "slide_name": slide_layout_name,
                "contents": {
                    title_shape_name: {
                        "text": "this is test title",
                    },
                    body_shape_name: {
                        "text": "test body",
                    }
                }
            }
        ]
        make_ppt(self.master_slide, self.target_slide, input)
        ppt = self.get_target_ppt()
        print(ppt.slides)
        self.assertEqual(1, len(ppt.slides))
        shapes = ppt.slides[0].shapes
        self.assertEqual(2, len(shapes))

        slide = ppt.slides[0]
        contents = input[0]['contents']
        title_shape = self.find_shape(slide, title_shape_name)
        self.assertEqual(contents[title_shape_name]['text'], title_shape.text)

        body_shape = self.find_shape(slide, body_shape_name)
        self.assertEqual(contents[body_shape_name]['text'], body_shape.text)

    def test_dirty_shape_name(self):
        ms = self.get_master_ppt()
        slide_layout_name = 'dirty_shape_name'
        origin_layout = ms.slide_master.slide_layouts.get_by_name(slide_layout_name)
        # check placeholder exists
        title_shape_name = "TITLE@@ AB"
        body_shape_name = "This is Body"
        self.assertTrue([s for s in origin_layout.shapes if s.name == title_shape_name][0])
        self.assertTrue([s for s in origin_layout.shapes if s.name == body_shape_name][0])

        input = [
            {
                "slide_name": slide_layout_name,
                "contents": {
                    title_shape_name: {
                        "text": "this is test title",
                    },
                    body_shape_name: {
                        "text": "test body",
                    }
                }
            }
        ]
        make_ppt(self.master_slide, self.target_slide, input)
        ppt = self.get_target_ppt()
        print(ppt.slides)
        self.assertEqual(1, len(ppt.slides))
        shapes = ppt.slides[0].shapes
        self.assertEqual(2, len(shapes))

        slide = ppt.slides[0]
        contents = input[0]['contents']
        title_shape = self.find_shape(slide, title_shape_name)
        self.assertEqual(contents[title_shape_name]['text'], title_shape.text)

        body_shape = self.find_shape(slide, body_shape_name)
        self.assertEqual(contents[body_shape_name]['text'], body_shape.text)
