from typing import Dict, List, Union

from pptx import Presentation

from .replace import replace_data
from .type import ReplaceSlideType, SlideData


def delete_slide(prs, slide):
    # copy from https://github.com/scanny/python-pptx/issues/67#issuecomment-568921262
    id_dict = {slide.id: [i, slide.rId] for i, slide in enumerate(prs.slides._sldIdLst)}
    slide_id = slide.slide_id
    prs.part.drop_rel(id_dict[slide_id][1])
    del prs.slides._sldIdLst[id_dict[slide_id][0]]


def get_slide_by_pos(prs, pos: int):
    return prs.slides[int(pos) - 1]


def make_ppt(
        master_file: str, target_file: str,
        pages: List[Union[Dict, ReplaceSlideType, SlideData]] = None, delete_pages: List[int] = None,
        mode: str = "replace"):
    master_pt = Presentation(master_file)
    if pages:
        for page in pages:
            if mode == 'template':
                raise NotImplementedError("doesn't support Template mode")
            elif mode == 'replace':
                if isinstance(page, dict):
                    page = ReplaceSlideType.from_dict(page)
                slide = get_slide_by_pos(master_pt, page.slide_pos)
            else:
                raise NotImplementedError(f"doesn't support {mode} mode")
            replace_data(slide, page.contents)
    if delete_pages:
        delete_slides = [get_slide_by_pos(master_pt, pos) for pos in delete_pages]
        for slide in delete_slides:
            delete_slide(master_pt, slide)
        print(f'delete {len(delete_pages)} slide')

    master_pt.save(target_file)
