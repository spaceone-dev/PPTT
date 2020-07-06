from marshmallow import Schema, fields
from pptx.presentation import Presentation
from pptx.shapes.base import BaseShape
from pptx.slide import SlideMaster, SlideLayout


class ShapeSchema(Schema):
    name = fields.Str(required=True)
    is_placeholder = fields.Bool()
    type = fields.Str()

    @classmethod
    def load_by_pptx(cls, _sh: BaseShape):
        sh = cls()
        sh.name = _sh.name
        sh.is_placeholder = _sh.is_placeholder
        sh.type = ''
        if _sh.has_table:
            sh.type = 'table'
        elif _sh.has_chart:
            sh.type = 'chart'
        elif _sh.has_text_frame:
            sh.type = 'text'
        return sh


class SlideLayoutSchema(Schema):
    shape_schema_class = ShapeSchema
    name = fields.Str(required=True)
    shapes = fields.List(fields.Nested(ShapeSchema))

    @classmethod
    def load_by_pptx(cls, _sl: SlideLayout):
        sl = cls()
        sl.name = _sl.name
        sl.shapes = [cls.shape_schema_class.load_by_pptx(sh) for sh in _sl.shapes]
        return sl


class SlideMasterSchema(Schema):
    slide_layout_class = SlideLayoutSchema
    name = fields.Str(required=True)
    main = fields.Bool(default=True)
    slides = fields.List(fields.Nested(SlideLayoutSchema))

    @classmethod
    def load_by_pptx(cls, _sm: SlideMaster, main):
        sm = cls()
        sm.name = _sm.name
        sm.main = main
        sm.slides = [cls.slide_layout_class.load_by_pptx(sl) for sl in _sm.slide_layouts]
        return sm


class PPTTemplate(Schema):
    slide_master_class = SlideMasterSchema
    main = fields.Str()
    masters = fields.List(fields.Nested(SlideMasterSchema))

    @classmethod
    def load_by_pptx(cls, _ppt: Presentation):
        ppt = cls()
        ppt.name = _ppt.slide_master.name
        ppt.masters = [cls.slide_master_class.load_by_pptx(sm, idx == 0) for idx, sm in enumerate(_ppt.slide_masters)]
        return ppt


class SlideStub(Schema):
    _slide_layout_name = None

    def get_contents(self):
        result = {}
        for name, field_type in self.fields.items():
            if isinstance(field_type, fields.String):
                try:
                    value = getattr(self, name)
                    result[name] = {"text": value}
                except Exception as e:
                    print(e)

        return result
