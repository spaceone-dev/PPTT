from ast import Module, ImportFrom, alias, Assign, Name, Store, List, Load, Subscript, Index, ClassDef, Constant, \
    AnnAssign, keyword, Call

from ast_decompiler import decompile
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.shapetree import SlideShapes
from pptx.slide import Slides, Slide

from ..utils import name_to_slugify

dataclass_decorator = Name(id='dataclass', ctx=Load())


def ast_index(name: str):
    return Index(value=Name(id=name, ctx=Load()))


def ast_optional_subscript(name: str):
    """
    same as this
    Subscript(
        value=Name(id='Optional', ctx=Load()),
        slice=Index(value=Name(id= <name> , ctx=Load())),
        ctx=Load()
    )
    """
    return Subscript(
        value=Name(id='Optional', ctx=Load()),
        slice=ast_index(name),
        ctx=Load()
    )


def ast_default_field(value):
    return Call(
        func=Name(id='field', ctx=Load()),
        args=[],
        keywords=[keyword(arg='default', value=Constant(value=value, kind=None))]
    )


text_data_value = ast_optional_subscript('TextData')

chart_xy_data_value = ast_optional_subscript('ChartXYData')
chart_category_data_value = ast_optional_subscript('ChartCategoryData')
chart_bubble_data_value = ast_optional_subscript('ChartBubbleData')
table_data_value = ast_optional_subscript('TableData')

imports = [
    ImportFrom(
        module='dataclasses',
        names=[
            alias(name='dataclass', asname=None),
            alias(name='field', asname=None),
        ],
        level=0
    ),
    ImportFrom(module='typing', names=[alias(name='Optional', asname=None)], level=0),
    ImportFrom(
        module='PPTT.type',
        names=[
            alias(name='SlideData', asname=None),
            alias(name='TextData', asname=None),
            alias(name='ChartCategoryData', asname=None),
            alias(name='ChartBubbleData', asname=None),
            alias(name='ChartXYData', asname=None),
            alias(name='TableData', asname=None),
        ],
        level=0),
]


def make_slides_assign(slides_count):
    return Assign(
        targets=[Name(id='SLIDES', ctx=Store())],
        value=List(
            elts=[Name(id=f'Slide{idx}', ctx=Load()) for idx in range(1, slides_count + 1)],
            ctx=Load()
        ),
        type_comment=None
    )


CHART_DATA_TYPE_MAP = {
    XL_CHART_TYPE.BUBBLE: chart_bubble_data_value,
    XL_CHART_TYPE.BUBBLE_THREE_D_EFFECT: chart_bubble_data_value,
    XL_CHART_TYPE.XY_SCATTER: chart_xy_data_value,
    XL_CHART_TYPE.XY_SCATTER_LINES: chart_xy_data_value,
    XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS: chart_xy_data_value,
    XL_CHART_TYPE.XY_SCATTER_SMOOTH: chart_xy_data_value,
    XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS: chart_xy_data_value,
}


def make_data_field(shape: GraphicFrame):
    if shape.has_text_frame:
        value = text_data_value
    elif shape.has_table:
        value = table_data_value
    elif shape.has_chart:
        value = CHART_DATA_TYPE_MAP.get(shape.chart.chart_type, chart_category_data_value)
    else:
        return None

    return AnnAssign(
        target=Name(id=name_to_slugify(shape.name), ctx=Store()),
        annotation=value,
        value=ast_default_field(None),
        simple=1,
    )


def make_slide_content_class(name: str, shapes: SlideShapes):
    body = []
    for shape in shapes:
        if field := make_data_field(shape):
            body.append(field)
    return ClassDef(
        name=name,
        bases=[], keywords=[],
        body=body,
        decorator_list=[dataclass_decorator]
    )


def make_slide_class(slide: Slide, slide_index: int) -> list:
    slide_class_name = f"Slide{slide_index}"
    content_class_name = f"{slide_class_name}Content"
    return [
        make_slide_content_class(content_class_name, slide.shapes),
        ClassDef(
            name=slide_class_name,
            bases=[Name(id='SlideData', ctx=Load())],
            keywords=[],
            body=[
                AnnAssign(
                    target=Name(id='contents', ctx=Store()),
                    annotation=ast_optional_subscript(content_class_name),
                    value=ast_default_field(None),
                    simple=1
                ),
                AnnAssign(
                    target=Name(id='slide_pos', ctx=Store()),
                    annotation=Name(id='int', ctx=Load()),
                    value=Constant(value=slide_index, kind=None),
                    simple=1
                )
            ],
            decorator_list=[dataclass_decorator]),
    ]


def generate_ast(slides: Slides) -> Module:
    classes = []
    for idx, slide in enumerate(slides):
        classes += make_slide_class(slide, idx + 1)
    body = imports + classes + [make_slides_assign(len(slides))]
    return Module(body=body, type_ignores=[])


def make_stub(pptx_path: str, output_path: str):
    master_ppt = Presentation(pptx_path)
    ast_result = generate_ast(master_ppt.slides)
    with open(output_path, mode='w') as fp:
        fp.write(decompile(ast_result))
