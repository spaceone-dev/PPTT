import logging
from dataclasses import dataclass, is_dataclass
from itertools import repeat
from typing import List, Callable, Dict, Union

from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
from pptx.compat import to_unicode
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.shapes.graphfrm import GraphicFrame
from pptx.text.text import TextFrame, _Run, _Paragraph

from ..type import KeyValueData, CategoryData, KVKeys, ChartData, ChartDataTypeHints, XYData, \
    BubbleData, RawData, TextFrameDataType, TextStyle, KVKey, SlideContentsType, TextData, TableData
from ..utils import find_shape, find_shape_by_slide_layout


def replace_paragraphs_text(paragraph, new_text):
    # copy from https://github.com/scanny/python-pptx/issues/285#issuecomment-300273686
    p = paragraph._p  # the lxml element containing the `<a:p>` paragraph element
    # remove all except first run
    if paragraph.runs:
        for idx, run in enumerate(paragraph.runs):
            if idx == 0:
                continue
            p.remove(run._r)
    else:
        paragraph.add_run()
    paragraph.runs[0].text = new_text


def copy_font_style(style_run, target_run):
    target_run.font.name = style_run.font.name
    target_run.font.size = style_run.font.size
    target_run.font.bold = style_run.font.bold
    target_run.font.italic = style_run.font.italic
    target_run.font.underline = style_run.font.underline
    target_run.font.language_id = style_run.font.language_id

    if color_type := style_run.font.color.type:
        if color_type == MSO_COLOR_TYPE.RGB:
            target_run.font.color.rgb = style_run.font.color.rgb
        elif color_type == MSO_COLOR_TYPE.SCHEME:
            target_run.font.color.theme_color = style_run.font.color.theme_color
            target_run.font.color.brightness = style_run.font.color.brightness


def copy_paragraph_style(style_paragraph, target_paragraph):
    style_paragraph.alignment = target_paragraph.alignment
    style_paragraph.level = target_paragraph.level
    style_paragraph.line_spacing = target_paragraph.line_spacing
    style_paragraph.space_after = target_paragraph.space_after
    style_paragraph.space_before = target_paragraph.space_before
    copy_font_style(style_paragraph, target_paragraph)


def clear_text_frame(text_frame: TextFrame, all: bool = False):
    index = 0 if all else 1
    # remove all paragraphs
    for p in text_frame._txBody.p_lst[index:]:
        text_frame._txBody.remove(p)


def replace_font_style(subshape: Union[_Run, _Paragraph], text_frame_data: TextStyle):
    if font_style := text_frame_data.font:
        if font_color := font_style.color:
            # font_color value must be hex cololr ex) #ffffff
            subshape.font.color.rgb = RGBColor.from_string(font_color[1:])
        if isinstance(font_style.bold, bool):
            subshape.font.bold = font_style.bold
        if isinstance(font_style.italic, bool):
            subshape.font.italic = font_style.italic
        if isinstance(font_style.underline, bool):
            subshape.font.underline = font_style.underline


def replace_text_frame(text_frame: TextFrame, text_frame_data: TextFrameDataType):
    clear_text_frame(text_frame)
    new_text = text_frame_data.value if isinstance(text_frame_data, TextStyle) else text_frame_data
    if new_text is None:
        new_text = ''

    style_paragraph = text_frame.paragraphs[0]
    style_run = style_paragraph.runs[0] if len(style_paragraph.runs) else style_paragraph.add_run()
    if isinstance(text_frame_data, TextStyle):  # if value type is TextStyleType
        replace_font_style(style_paragraph, text_frame_data)
        replace_font_style(style_run, text_frame_data)

    split_text = to_unicode(f"{new_text}").split('\n')
    for line_num, line_text in enumerate(split_text):
        paragraph = text_frame.add_paragraph() if line_num != 0 else text_frame.paragraphs[0]
        replace_paragraphs_text(paragraph, line_text)
        copy_paragraph_style(style_paragraph, paragraph)
        copy_font_style(style_run, paragraph.runs[0])


def replace_table_data_raw(shape: GraphicFrame, data: RawData):
    records = data.data or []
    for r_idx, row in enumerate(shape.table.rows):
        try:
            record = records[r_idx]
        except IndexError:
            record = list(repeat(None, len(row.cells)))
        for c_idx, cell in enumerate(row.cells):
            try:
                value = record[c_idx]
            except IndexError:
                value = ''
            replace_text_frame(cell.text_frame, value)


def replace_table_data_key_value(shape: GraphicFrame, data: KeyValueData):
    keys: KVKeys = data.keys or []
    records = data.data or []
    header_names: List[str] = [k.name or '' if isinstance(k, KVKey) else k for k in keys]
    data_keys: List[str] = [k.data_key or '' if isinstance(k, KVKey) else k for k in keys]

    raw_records = [[record.get(key) for key in data_keys] for record in records]
    raw_data = [header_names] + raw_records
    replace_table_data_raw(shape, RawData(data=raw_data))


TABLE_DATA_TYPE_HANDLER = {
    "key_value": replace_table_data_key_value,
    "raw": replace_table_data_raw,
}


def replace_category_data(shape: GraphicFrame, data: CategoryData):
    chart_data = CategoryChartData()
    chart_data.categories = data.categories
    for name, series_data in data.series.items():
        chart_data.add_series(name, series_data)
    shape.chart.replace_data(chart_data)


def replace_xy_data(shape: GraphicFrame, data: XYData):
    chart_data = XyChartData()

    for name, points in data.series.items():
        series_data = chart_data.add_series(name)
        for point in points:
            series_data.add_data_point(point.x, point.y)

    shape.chart.replace_data(chart_data)


def replace_bubble_data(shape: GraphicFrame, data: BubbleData):
    chart_data = BubbleChartData()

    for name, points in data.series.items():
        series_data = chart_data.add_series(name)
        for point in points:
            series_data.add_data_point(point.x, point.y, point.size)

    shape.chart.replace_data(chart_data)


CHART_DATA_TYPE_HANDLER: Dict[str, Callable[[GraphicFrame, ChartDataTypeHints], None]] = {
    "category_data": replace_category_data,
    "xy_data": replace_xy_data,
    "bubble_data": replace_bubble_data,
}


def replace_chart_data(shape: GraphicFrame, chart_data: ChartDataTypeHints):
    CHART_DATA_TYPE_HANDLER[chart_data.data_type](shape, chart_data)
    if chart_data.title:
        replace_text_frame(shape.chart.chart_title.text_frame, chart_data.title)


def replace_data(slide, contents: Union[SlideContentsType, dataclass], mode="replace"):
    _contents = []

    if isinstance(contents, dict):
        _contents = contents.items()
    elif is_dataclass(contents):
        _contents = ((name, getattr(contents, name, None)) for name in contents.__dataclass_fields__.keys())
    else:
        raise ValueError('un support type')
    for name, data in _contents:
        shape = find_shape(slide, name) if mode == 'replace' else find_shape_by_slide_layout(slide, name)
        if shape and data:
            # shape.name = name # for
            if isinstance(data, TextData):
                replace_text_frame(shape.text_frame, data.text)
            elif isinstance(data, TableData):
                try:
                    TABLE_DATA_TYPE_HANDLER[data.table.data_type](shape, data.table)
                except Exception as e:
                    logging.exception(e)

            elif isinstance(data, ChartData):
                try:
                    replace_chart_data(shape, data.chart)
                except Exception as e:
                    logging.exception(e)
